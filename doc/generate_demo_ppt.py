#!/usr/bin/env python3
"""Generate wav2chat product demo PowerPoint."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "wav2chat-demo.pptx"

# WeChat-inspired accent + neutral slides
COLOR_TITLE = RGBColor(0x1A, 0x1A, 0x2E)
COLOR_ACCENT = RGBColor(0x07, 0xC1, 0x60)
COLOR_SUB = RGBColor(0x55, 0x55, 0x55)
COLOR_LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_BUBBLE_ME = RGBColor(0x95, 0xEC, 0x69)
COLOR_BUBBLE_OTHER = RGBColor(0xFF, 0xFF, 0xFF)


def _set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, COLOR_TITLE)

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.alignment = PP_ALIGN.LEFT

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(1))
    stf = sub.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(22)
    sp.font.color.rgb = COLOR_BUBBLE_ME
    sp.alignment = PP_ALIGN.LEFT

    tag = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11), Inches(0.5))
    ttf = tag.text_frame
    ttf.text = "Python · FunASR · ffmpeg · wxPython  |  CLI + 桌面 GUI  |  AGPL-3.0"
    tp = ttf.paragraphs[0]
    tp.font.size = Pt(14)
    tp.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)


def _add_section_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, COLOR_LIGHT_BG)

    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    htf = header.text_frame
    htf.text = title
    hp = htf.paragraphs[0]
    hp.font.size = Pt(32)
    hp.font.bold = True
    hp.font.color.rgb = COLOR_TITLE

    accent = slide.shapes.add_shape(
        1, Inches(0.6), Inches(1.15), Inches(1.2), Inches(0.08)
    )  # rectangle
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_ACCENT
    accent.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(5.5))
    btf = body.text_frame
    btf.word_wrap = True
    for i, line in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_SUB
        p.space_after = Pt(10)


def _add_code_slide(prs: Presentation, title: str, code: str, note: str = "") -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, COLOR_LIGHT_BG)

    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    htf = header.text_frame
    htf.text = title
    hp = htf.paragraphs[0]
    hp.font.size = Pt(28)
    hp.font.bold = True
    hp.font.color.rgb = COLOR_TITLE

    panel = slide.shapes.add_shape(1, Inches(0.7), Inches(1.3), Inches(11.9), Inches(4.8))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(0x2D, 0x2D, 0x2D)
    panel.line.color.rgb = RGBColor(0x44, 0x44, 0x44)

    code_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.3), Inches(4.4))
    ctf = code_box.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(code.strip().split("\n")):
        p = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)

    if note:
        nb = slide.shapes.add_textbox(Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.8))
        ntf = nb.text_frame
        ntf.text = note
        np = ntf.paragraphs[0]
        np.font.size = Pt(16)
        np.font.color.rgb = COLOR_SUB


def _add_flow_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, COLOR_LIGHT_BG)

    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    header.text_frame.text = "处理流程"
    header.text_frame.paragraphs[0].font.size = Pt(32)
    header.text_frame.paragraphs[0].font.bold = True
    header.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    steps = [
        ("音频文件", "wav / mp3 / m4a …"),
        ("ffmpeg", "mono · 16 kHz"),
        ("FunASR", "VAD→ASR→标点→说话人"),
        ("输出", "txt + json"),
    ]
    x = 0.5
    for i, (name, desc) in enumerate(steps):
        box = slide.shapes.add_shape(1, Inches(x), Inches(2.5), Inches(2.6), Inches(1.6))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        box.line.color.rgb = COLOR_ACCENT
        tf = box.text_frame
        tf.text = name
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_TITLE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = COLOR_SUB
        p2.alignment = PP_ALIGN.CENTER
        if i < len(steps) - 1:
            arrow = slide.shapes.add_textbox(
                Inches(x + 2.65), Inches(3.0), Inches(0.4), Inches(0.5)
            )
            arrow.text_frame.text = "→"
            arrow.text_frame.paragraphs[0].font.size = Pt(28)
            arrow.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        x += 3.15

    models = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(11.8), Inches(1.8))
    mtf = models.text_frame
    mtf.text = (
        "FunASR 四步：fsmn-vad → paraformer-zh → ct-punc → cam++  |  "
        "非单次 ASR，CPU 上长录音较慢属正常"
    )
    mtf.paragraphs[0].font.size = Pt(16)
    mtf.paragraphs[0].font.color.rgb = COLOR_SUB


def _add_gui_mock_slide(prs: Presentation) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_slide_bg(slide, COLOR_LIGHT_BG)

    header = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    header.text_frame.text = "桌面 GUI — 像读微信一样读录音"
    header.text_frame.paragraphs[0].font.size = Pt(28)
    header.text_frame.paragraphs[0].font.bold = True
    header.text_frame.paragraphs[0].font.color.rgb = COLOR_TITLE

    # Left panel mock
    left = slide.shapes.add_shape(1, Inches(0.6), Inches(1.4), Inches(3.2), Inches(5.2))
    left.fill.solid()
    left.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    left.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    ltf = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(2.8), Inches(4.8))
    ltf.text_frame.text = (
        "文件队列\n\n● call.m4a  ✓\n● meeting.wav\n\n"
        "人数 [2] 至 [2] 人\n☐ 刷新模型\n\n[转换] ☑ 自动转换\n\n拖放 · Delete 删除"
    )
    for p in ltf.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_SUB

    # Right chat mock
    right = slide.shapes.add_shape(1, Inches(4.1), Inches(1.4), Inches(8.4), Inches(5.2))
    right.fill.solid()
    right.fill.fore_color.rgb = RGBColor(0xEC, 0xEC, 0xEC)
    right.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)

    # Other bubble (left)
    b1 = slide.shapes.add_shape(1, Inches(4.5), Inches(2.0), Inches(3.5), Inches(0.7))
    b1.fill.solid()
    b1.fill.fore_color.rgb = COLOR_BUBBLE_OTHER
    b1.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    b1.text_frame.text = "喂，你好。"
    b1.text_frame.paragraphs[0].font.size = Pt(14)

    # Me bubble (right)
    b2 = slide.shapes.add_shape(1, Inches(8.2), Inches(3.0), Inches(3.8), Inches(0.7))
    b2.fill.solid()
    b2.fill.fore_color.rgb = COLOR_BUBBLE_ME
    b2.line.fill.background()
    b2.text_frame.text = "你好，我想问一下续贷的事情。"
    b2.text_frame.paragraphs[0].font.size = Pt(14)
    b2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

    feats = slide.shapes.add_textbox(Inches(0.7), Inches(6.75), Inches(11.8), Inches(0.6))
    feats.text_frame.text = (
        "列表 / 气泡视图 · 拖放批量导入 · emoji 头像 · 右键设「这是我」· 说话人资料编辑 · 搜索"
    )
    feats.text_frame.paragraphs[0].font.size = Pt(14)
    feats.text_frame.paragraphs[0].font.color.rgb = COLOR_SUB


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(
        prs,
        "wav2chat",
        "把电话 / 会议录音转换成按说话人分段的聊天文本",
    )

    _add_section_slide(
        prs,
        "为什么需要 wav2chat？",
        [
            "电话、会议录音往往长达数十分钟，逐段回听效率极低",
            "通用语音识别只给一大段文字，分不清谁说了什么",
            "纯文本 transcript 缺少对话结构，复盘和归档都不直观",
            "wav2chat：本地运行 · 自动说话人分离 · 聊天式阅读体验",
        ],
    )

    _add_section_slide(
        prs,
        "核心功能",
        [
            "支持 wav / mp3 / m4a / amr / aac / flac / ogg 等常见格式",
            "自动 VAD、中文标点恢复、说话人分离（FunASR cam++）",
            "输出同名 .txt（可读对话）与 .json（结构化数据，可再次打开编辑）",
            "命令行批处理整个目录；GUI 拖放队列 + 一键转换",
            "说话人显示名、emoji 头像、主角（右侧绿色气泡）可自定义",
            "面向中文电话 / 会议场景优化；无需数据库与联网服务（模型首次需下载）",
        ],
    )

    _add_flow_slide(prs)

    _add_code_slide(
        prs,
        "命令行 — 快速上手",
        """# 转换单个文件（生成 call.txt）
wav2chat call.m4a

# 同时输出 JSON，双人电话场景
wav2chat call.wav -o call.txt --json call.json \\
  --min-speakers 2 --max-speakers 2 \\
  --role spk0=对方 --role spk1=我

# 批量处理目录
wav2chat ./recordings --batch -o ./texts --json ./jsons

# 打开图形界面
wav2chat -g""",
        "已有旧版 JSON？运行 jsonfix *.json 迁移到新 speakers 格式",
    )

    _add_gui_mock_slide(prs)

    _add_section_slide(
        prs,
        "GUI 亮点",
        [
            "拖放文件 / 文件夹，导入进度对话框，已有 JSON 则直接加载、不重复转写",
            "人数 [min] 至 [max]（默认 2–2）；☐ 刷新模型（cache 异常时用）",
            "文件队列：未转换 / 转换中 / 已完成；Delete 删除选中项",
            "气泡视图：仿微信布局，「我」在右侧绿色气泡，对方在左侧",
            "头像右键：设为主角、编辑资料；列表视图 + 搜索；多语言界面",
            "状态栏 / 日志：加载模型 → 规范化 → 语音识别 → 保存",
        ],
    )

    _add_code_slide(
        prs,
        "JSON 输出结构（新版）",
        """{
  "source": "call.m4a",
  "duration": 126.5,
  "primary_speaker": 1,
  "speakers": [
    { "name": "spk0", "role": "对方", "gender": "", "avatar": "👦" },
    { "name": "spk1", "role": "me",   "gender": "", "avatar": "👧" }
  ],
  "segments": [
    { "start": 1.2, "end": 3.8, "speaker": 0, "text": "喂，你好。" },
    { "start": 4.1, "end": 7.6, "speaker": 1, "text": "你好，我想问一下…" }
  ]
}""",
        "spk1 默认 role 为 me；可用 jsonfix 批量补全/迁移历史文件",
    )

    _add_section_slide(
        prs,
        "为什么转换需要时间？",
        [
            "加载模型：4 个神经网络首次读入内存，同一会话约 10–40 s，之后跳过",
            "规范化：ffmpeg → mono 16 kHz，通常数秒",
            "转写：VAD + 识别 + 标点 + 说话人分离，CPU 常见 0.2–1.0 倍实时",
            "例：10 分钟电话 → CPU 约 2–10 分钟；GPU 可快一个数量级",
            "CPU 上 FunASR 基本逐 VAD 段推理，长录音必然偏慢",
            "GUI 进度百分比含心跳估算，开日志看「语音识别」阶段更准确",
        ],
    )

    _add_section_slide(
        prs,
        "模型加载优化（cache / mmap）",
        [
            "首次：ModelScope 下载到 ~/.cache/modelscope/（一次性）",
            "默认：本地 cache 路径 + torch.load(mmap) + 跳过 hub 更新检查",
            "同一会话 GUI/CLI 模型常驻内存，第二次转换不再读盘",
            "jieba（标点）：~/.cache/wav2chat/jieba/ 持久 cache，词典变更才 rebuild",
            "异常时勾选「刷新模型」或 wav2chat --refresh-models 强制重载",
            "环境检查：wav2chat --version",
        ],
    )

    _add_section_slide(
        prs,
        "速度与质量 — 可调参数",
        [
            "【已暴露】GPU / CPU：装 CUDA 版 PyTorch，收益最大",
            "【已暴露】人数 min–max：电话默认 2–2，收窄范围加快聚类、提高准确率",
            "【已暴露】--refresh-models：仅 cache 损坏或需更新时使用",
            "【已暴露】-v / 日志面板：查看 Transcribing、RTF、各 ckpt 耗时",
            "【内置默认】batch_size_s=300；batch_size_threshold_s=60（FunASR）",
            "【内置】四模型全开（含 cam++ 分离）；跳过分离可更快，当前未提供快捷开关",
        ],
    )

    _add_code_slide(
        prs,
        "调优命令示例",
        """# 双人电话（GUI 默认 2–2）
wav2chat call.wav -n 2 -m 2 -o call.txt --json call.json

# 强制从 ModelScope 刷新模型
wav2chat --refresh-models call.m4a

# 查看 torch / funasr / ffmpeg 环境
wav2chat --version

# 长任务：开 verbose 看真实进度
wav2chat call.m4a -v -n 2 -m 2""",
        "实用建议：优先 GPU · 保持 GUI 不关 · 勿滥用刷新模型 · 说话人不对调 --role",
    )

    _add_section_slide(
        prs,
        "安装与运行环境",
        [
            "系统：Linux（Debian/Ubuntu 推荐）+ ffmpeg",
            "Python 3.10+ · funasr · modelscope · torch（可选 GPU）",
            "GUI：系统包 python3-wxgtk4.0 + pip install -e .",
            "首次运行会从 ModelScope 下载模型，请预留磁盘与时间",
            "安装：pip install -e .   →   wav2chat --help   →   wav2chat -g",
        ],
    )

    _add_section_slide(
        prs,
        "总结",
        [
            "wav2chat = 录音 → 分段对话文本 → 聊天式阅读",
            "CLI 批处理 + GUI 拖放/气泡/头像；本地运行，无 Web 服务",
            "调优：GPU · 2–2 人数 · cache+mmap · 会话内复用模型",
            "开源：GNU AGPL v3  |  详见 README-zh.md 性能与调优",
            "",
            "谢谢！欢迎试用与反馈",
        ],
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
